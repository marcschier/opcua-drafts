"""Structural checks on the generated deck.

There is no headless renderer in this repository, so overflow is estimated rather
than observed: for every text-bearing shape the checker measures how much text was
asked to fit and compares it against the space available.

    python decks/check_layout.py                  # check the default deck
    python decks/check_layout.py --deck other.pptx
    python decks/check_layout.py --strict         # non-zero exit on any finding

Findings are advisory. A slide flagged here is worth opening; a slide not flagged
is not guaranteed perfect.
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

import build_deck
import theme
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

DEFAULT_DECK = Path(__file__).parent / "OPC-UA-Drafts-Overview.pptx"

# Segoe UI averages a little under 0.5 em per character across mixed-case prose.
AVG_CHAR_EM = 0.50
LINE_SPACING = 1.18
EMU_PER_PT = 12700
EMU_PER_INCH = 914400
OVERLAP_TOLERANCE = Emu(int(0.03 * EMU_PER_INCH))

TITLE_COLLISION_LAYOUTS = {"bullets", "table", "code", "diagram", "demo"}


def text(shape) -> str:
    """
    Return the shape's text, or an empty string for non-text shapes.
    """
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def rect(shape) -> tuple[int, int, int, int]:
    """
    Return ``(left, top, right, bottom)`` in EMUs.
    """
    return (
        int(shape.left),
        int(shape.top),
        int(shape.left + shape.width),
        int(shape.top + shape.height),
    )


def shape_label(shape) -> str:
    """
    Human-readable shape label for findings.
    """
    value = text(shape)
    if value:
        return value.replace("\n", " ").encode("ascii", "replace").decode("ascii")[:50]
    return str(shape.shape_type)


def slide_specs() -> list[dict]:
    """
    Load source slide metadata so layout-sensitive geometry checks know intent.
    """
    docs, _ = build_deck.load_content()
    specs: list[dict] = []
    for doc in docs:
        specs.extend(doc.get("slides", []) or [])
    return specs


def is_full_bleed_background(shape, slide_area: int) -> bool:
    """
    Title, section and statement slides use zero-text panels behind all content.
    """
    if text(shape):
        return False
    area = int(shape.width) * int(shape.height)
    return area > slide_area * 0.70


def is_accent_bar(shape, width: int, height: int) -> bool:
    """
    Accent bars are decoration: thin, zero-text strips on the top or left edge.
    """
    if text(shape):
        return False
    thin_h = int(shape.height) <= int(Inches(0.20))
    thin_w = int(shape.width) <= int(Inches(0.20))
    spans_width = int(shape.width) >= width * 0.50
    spans_height = int(shape.height) >= height * 0.50
    return (thin_h and spans_width) or (thin_w and spans_height)


def is_chrome_text(shape) -> bool:
    """
    Shared title, kicker, footer text and page number are not slide content.
    """
    if not text(shape):
        return False
    top = int(shape.top)
    left = int(shape.left)
    bottom = int(shape.top + shape.height)
    if top >= int(theme.FOOTER_TOP - Inches(0.12)):
        return True
    if bottom <= int(theme.BODY_TOP + Inches(0.04)):
        return True
    # Demo badges sit in the title band and are layout chrome, not body content.
    if top < int(theme.BODY_TOP) and left > int(theme.SLIDE_WIDTH * 0.70):
        return True
    return False


def is_decoration(shape, width: int, height: int) -> bool:
    """
    Non-content furniture should not participate in body geometry checks.
    """
    slide_area = width * height
    return (
        is_full_bleed_background(shape, slide_area)
        or is_accent_bar(shape, width, height)
        or is_chrome_text(shape)
    )


def is_connector(shape) -> bool:
    """
    Connectors intentionally have no content rectangle for these checks.
    """
    shape_type = str(shape.shape_type)
    return "LINE" in shape_type or "CONNECTOR" in shape_type


def is_arrow_label(shape, spec: dict | None) -> bool:
    """
    Diagram arrow labels are small text boxes deliberately centered on connectors.
    """
    if not spec or spec.get("layout") != "diagram" or not text(shape):
        return False
    labels = {str(arrow.get("label", "")).strip() for arrow in spec.get("arrows", []) or []}
    if text(shape) not in labels:
        return False
    return int(shape.width) <= int(Inches(1.75)) and int(shape.height) <= int(Inches(0.34))


def is_flow_text_region(shape) -> bool:
    """
    Bullet body text boxes reserve flow space; their actual glyphs may not fill it.
    """
    return (
        "TEXT_BOX" in str(shape.shape_type)
        and int(shape.width) >= int(theme.CONTENT_W * 0.80)
        and int(shape.height) >= int(Inches(3.0))
    )


def content_shapes(slide, width: int, height: int) -> list:
    """
    Return shapes with content geometry, excluding decorative slide furniture.
    """
    shapes = []
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        if is_connector(shape) or is_decoration(shape, width, height):
            continue
        shapes.append(shape)
    return shapes


def shape_text_height(shape) -> float:
    """
    Estimate, in points, the height the text in ``shape`` needs at its own font sizes.
    """
    frame = shape.text_frame
    width_pt = (
        shape.width - frame.margin_left - frame.margin_right
    ) / EMU_PER_PT
    if width_pt <= 0:
        return 0.0
    total = (frame.margin_top + frame.margin_bottom) / EMU_PER_PT
    for para in frame.paragraphs:
        text = "".join(run.text for run in para.runs)
        size = max(
            [run.font.size.pt for run in para.runs if run.font.size is not None] or [18.0]
        )
        indent = 0.0
        p_pr = para._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        if p_pr is not None and p_pr.get("marL"):
            indent = int(p_pr.get("marL")) / EMU_PER_PT
        usable = max(width_pt - indent, size)
        chars_per_line = max(usable / (size * AVG_CHAR_EM), 1.0)
        lines = max(math.ceil(len(text) / chars_per_line), 1)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else LINE_SPACING
        total += lines * size * max(spacing, 1.0)
        if para.space_before is not None:
            total += para.space_before.pt
        if para.space_after is not None:
            total += para.space_after.pt
    return total


A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

SITE_NAMES = {0: "top", 1: "left", 2: "bottom", 3: "right"}


def connector_endpoints(shape) -> tuple[tuple[int, int], tuple[int, int]] | None:
    """
    Return ``((start shape id, start site), (end shape id, end site))`` for a connector.

    ``python-pptx`` exposes no reader for the connection a connector was attached
    to, so the ids come straight off ``a:stCxn`` and ``a:endCxn``. A connector that
    floats free of any shape returns ``None``.
    """
    props = shape._element.find(f"{P_NS}nvCxnSpPr")
    if props is None:
        return None
    non_visual = props.find(f"{P_NS}cNvCxnSpPr")
    if non_visual is None:
        return None
    start = non_visual.find(f"{A_NS}stCxn")
    end = non_visual.find(f"{A_NS}endCxn")
    if start is None or end is None:
        return None
    return (
        (int(start.get("id")), int(start.get("idx"))),
        (int(end.get("id")), int(end.get("idx"))),
    )


def connector_is_elbow(shape) -> bool:
    """
    An elbow connector is routed as right angles by PowerPoint; a straight one is not.
    """
    geometry = shape._element.find(f".//{A_NS}prstGeom")
    if geometry is None:
        return False
    return str(geometry.get("prst", "")).startswith("bentConnector")


def site_faces(site: int, source, target) -> bool:
    """
    Is the edge an arrow leaves or enters pointing at the other shape?

    Leaving the right edge to reach a shape on the left is what draws a line back
    across the shape it just left, which is the defect this catches.
    """
    tolerance = Inches(0.05)
    if site == 0:
        return target.top + target.height <= source.top + tolerance
    if site == 2:
        return target.top >= source.top + source.height - tolerance
    if site == 1:
        return target.left + target.width <= source.left + tolerance
    if site == 3:
        return target.left >= source.left + source.width - tolerance
    return True


def check_connectors(slide, index: int, title: str) -> list[str]:
    """
    Report arrows that are not drawn as horizontal and vertical runs.

    Two ways an arrow goes wrong. A straight connector between shapes that do not
    share a centre line is a diagonal. And an edge that points away from the other
    shape sends the line back across its own box before it can turn around.
    """
    findings: list[str] = []
    by_id = {shape.shape_id: shape for shape in slide.shapes if not is_connector(shape)}

    for shape in slide.shapes:
        if not is_connector(shape):
            continue
        endpoints = connector_endpoints(shape)
        if endpoints is None:
            continue
        (start_id, start_site), (end_id, end_site) = endpoints
        source = by_id.get(start_id)
        target = by_id.get(end_id)
        if source is None or target is None:
            continue

        if not site_faces(start_site, source, target):
            findings.append(
                f"slide {index} ({title}): arrow from '{shape_label(source)}' to "
                f"'{shape_label(target)}' leaves the {SITE_NAMES[start_site]} edge, "
                f"which faces away from the target"
            )
        if not site_faces(end_site, target, source):
            findings.append(
                f"slide {index} ({title}): arrow from '{shape_label(source)}' to "
                f"'{shape_label(target)}' enters the {SITE_NAMES[end_site]} edge, "
                f"which faces away from the source"
            )

        if connector_is_elbow(shape):
            continue

        aligned_x = abs(
            (source.left + source.width / 2) - (target.left + target.width / 2)
        ) <= Inches(0.06)
        aligned_y = abs(
            (source.top + source.height / 2) - (target.top + target.height / 2)
        ) <= Inches(0.06)
        if not (aligned_x or aligned_y):
            findings.append(
                f"slide {index} ({title}): straight arrow from '{shape_label(source)}' "
                f"to '{shape_label(target)}' is diagonal \u2014 the shapes share no centre line"
            )
    return findings


def check(deck_path: Path) -> list[str]:
    """
    Walk every slide and report shapes that overflow, sit off-slide, or run empty.
    """
    prs = Presentation(deck_path)
    width = prs.slide_width
    height = prs.slide_height
    specs = slide_specs()
    findings: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        spec = specs[index - 1] if index <= len(specs) else None
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0][:60]
                break

        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            if not is_connector(shape) and (shape.width <= 0 or shape.height <= 0):                findings.append(
                    f"slide {index} ({title}): shape '{shape_label(shape)}' has empty geometry"
                )
            if (
                shape.left < Emu(-1000)
                or shape.top < Emu(-1000)
                or shape.left + shape.width > width + Emu(1000)
                or shape.top + shape.height > height + Emu(1000)
            ):
                findings.append(
                    f"slide {index} ({title}): shape '{shape.shape_type}' extends past the slide"
                )
            if is_connector(shape) or is_decoration(shape, int(width), int(height)):
                continue
            if shape.top + shape.height > theme.FOOTER_TOP:
                findings.append(
                    f"slide {index} ({title}): shape '{shape_label(shape)}' intrudes into the footer"
                )
            if (
                spec
                and spec.get("layout") in TITLE_COLLISION_LAYOUTS
                and shape.top < theme.BODY_TOP - OVERLAP_TOLERANCE
            ):
                findings.append(
                    f"slide {index} ({title}): shape '{shape_label(shape)}' sits above the body area"
                )
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text.strip():
                continue
            needed = shape_text_height(shape)
            available = shape.height / EMU_PER_PT
            if needed > available * 1.06:
                findings.append(
                    f"slide {index} ({title}): text needs ~{needed:.0f}pt in a "
                    f"{available:.0f}pt box \u2014 {int(needed / available * 100)}% full"
                )

        if spec and spec.get("layout") not in {"title", "section", "statement"}:
            shapes = content_shapes(slide, int(width), int(height))
            for left_index, left_shape in enumerate(shapes):
                if (
                    left_shape.has_table
                    or is_arrow_label(left_shape, spec)
                    or is_flow_text_region(left_shape)
                ):
                    continue
                left_l, left_t, left_r, left_b = rect(left_shape)
                for right_shape in shapes[left_index + 1 :]:
                    if (
                        right_shape.has_table
                        or is_arrow_label(right_shape, spec)
                        or is_flow_text_region(right_shape)
                    ):
                        continue
                    right_l, right_t, right_r, right_b = rect(right_shape)
                    overlap_w = min(left_r, right_r) - max(left_l, right_l)
                    overlap_h = min(left_b, right_b) - max(left_t, right_t)
                    if overlap_w > OVERLAP_TOLERANCE and overlap_h > OVERLAP_TOLERANCE:
                        findings.append(
                            f"slide {index} ({title}): shape '{shape_label(left_shape)}' overlaps "
                            f"'{shape_label(right_shape)}'"
                        )

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes and len(notes) < 120:
                findings.append(f"slide {index} ({title}): speaker notes are very short")

        findings.extend(check_connectors(slide, index, title))

    for index, slide in enumerate(prs.slides, start=1):
        tables = [s for s in slide.shapes if s.has_table]
        for shape in tables:
            table = shape.table
            for r_index, row in enumerate(table.rows):
                for c_index, cell in enumerate(row.cells):
                    text = cell.text_frame.text
                    col_w = table.columns[c_index].width / EMU_PER_PT
                    size = 11.0
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                size = run.font.size.pt
                    chars_per_line = max(col_w / (size * AVG_CHAR_EM), 1.0)
                    lines = math.ceil(len(text) / chars_per_line)
                    if lines > 2:
                        findings.append(
                            f"slide {index}: table cell r{r_index}c{c_index} wraps to "
                            f"~{lines} lines \u2014 shorten it"
                        )
    return findings


def main(argv: list[str] | None = None) -> int:
    """
    Command-line entry point.
    """
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--strict", action="store_true")
    args = parser.parse_args(argv)

    if not args.deck.exists():
        print(f"{args.deck} does not exist \u2014 run build_deck.py first")
        return 1

    findings = check(args.deck)
    for finding in findings:
        print(finding)

    prs = Presentation(args.deck)
    print(f"\n{len(prs.slides)} slides, {len(findings)} findings")
    return 1 if (args.strict and findings) else 0


if __name__ == "__main__":
    raise SystemExit(main())
