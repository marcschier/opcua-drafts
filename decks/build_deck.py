"""Build the OPC UA drafts overview deck from the YAML under ``content/``.

The ``.pptx`` is a build artifact. Edit the YAML, never the PowerPoint file.

    python decks/build_deck.py                 # regenerate the deck
    python decks/build_deck.py --check         # validate content, write nothing
    python decks/build_deck.py --out other.pptx

Content model
-------------
One YAML file per effort (and one per demo) under ``content/``. Each file is::

    id: core-01-data-channels
    section: core                 # core | cloud | metaverse | wot | companion | front | close
    order: 110                    # global sort key
    footer: OPC UA - Data Channels
    slides:
      - layout: bullets
        title: Why it exists
        ...

Supported layouts: ``title``, ``section``, ``bullets``, ``two-column``, ``table``,
``code``, ``diagram``, ``statement``, ``demo``. See ``README.md`` for the field
reference.
"""

from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import theme

CONTENT_DIR = Path(__file__).parent / "content"
DEFAULT_OUT = Path(__file__).parent / "OPC-UA-Drafts-Overview.pptx"

LAYOUTS = {
    "title",
    "section",
    "bullets",
    "two-column",
    "table",
    "code",
    "diagram",
    "statement",
    "demo",
}


@dataclass
class Deck:
    """
    Accumulated build state: the presentation plus the problems found while building.
    """

    prs: Presentation
    problems: list[str] = field(default_factory=list)
    slide_count: int = 0

    def problem(self, where: str, message: str) -> None:
        """
        Record a content problem without aborting, so one run reports them all.
        """
        self.problems.append(f"{where}: {message}")


def load_content() -> tuple[list[dict[str, Any]], list[str]]:
    """
    Load every YAML file under ``content/`` and return them in presentation order.

    A file that fails to parse is reported rather than raised, so one malformed
    document does not hide the state of every other one.
    """
    files = sorted(CONTENT_DIR.rglob("*.yaml"))
    docs: list[dict[str, Any]] = []
    errors: list[str] = []
    for path in files:
        name = path.relative_to(CONTENT_DIR).as_posix()
        try:
            with path.open(encoding="utf-8") as handle:
                doc = yaml.safe_load(handle)
        except yaml.YAMLError as error:
            errors.append(f"{name}: not valid YAML \u2014 {error}")
            continue
        if not doc:
            continue
        if not isinstance(doc, dict):
            errors.append(f"{name}: top level must be a mapping")
            continue
        doc["_path"] = name
        docs.append(doc)
    docs.sort(key=lambda d: (d.get("order", 10_000), d["_path"]))
    return docs, errors


def new_presentation() -> Presentation:
    """
    Create an empty 16:9 presentation.
    """
    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT
    return prs


def add_slide(deck: Deck):
    """
    Append a blank slide.
    """
    blank = deck.prs.slide_layouts[6]
    deck.slide_count += 1
    return deck.prs.slides.add_slide(blank)


def textbox(slide, left, top, width, height):
    """
    Add a margin-free text box with word wrap enabled.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return box


def style_run(run, *, size, color, font=theme.FONT_BODY, bold=False, italic=False):
    """
    Apply the deck's typography to a single run.
    """
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def rich_runs(paragraph, text: str, *, size, color, bold=False, italic=False):
    """
    Emit runs for ``text``, rendering ``**bold**`` and `` `code` `` inline spans.
    """
    buffer: list[str] = []

    def flush() -> None:
        if not buffer:
            return
        run = paragraph.add_run()
        run.text = "".join(buffer)
        style_run(run, size=size, color=color, bold=bold, italic=italic)
        buffer.clear()

    index = 0
    while index < len(text):
        if text.startswith("**", index):
            end = text.find("**", index + 2)
            if end != -1:
                flush()
                run = paragraph.add_run()
                run.text = text[index + 2 : end]
                style_run(run, size=size, color=color, bold=True, italic=italic)
                index = end + 2
                continue
        if text[index] == "`":
            end = text.find("`", index + 1)
            if end != -1:
                flush()
                run = paragraph.add_run()
                run.text = text[index + 1 : end]
                style_run(
                    run,
                    size=Pt(size.pt - 1),
                    color=theme.BLUE_DEEP,
                    font=theme.FONT_MONO,
                    bold=bold,
                    italic=italic,
                )
                index = end + 1
                continue
        buffer.append(text[index])
        index += 1
    flush()


def set_indent(paragraph, left: int, hanging: int = 0) -> None:
    """
    Set the paragraph's left indent and hanging indent.

    ``python-pptx`` exposes no API for either, so the values are written straight
    onto ``a:pPr``. ``hanging`` is a positive amount pulled back from ``left``.
    """
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("marL", str(int(left)))
    p_pr.set("indent", str(-int(hanging)))


def set_notes(slide, notes: str | None) -> None:
    """
    Attach speaker notes to a slide.
    """
    if not notes:
        return
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes.strip()
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(11)
            run.font.name = theme.FONT_BODY


def chrome(deck: Deck, slide, spec: dict[str, Any], doc: dict[str, Any]) -> None:
    """
    Draw the shared slide furniture: accent bar, title, kicker and footer.
    """
    style = theme.section_style(doc.get("section", "front"))

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), theme.SLIDE_WIDTH, Inches(0.14)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.color
    bar.line.fill.background()
    bar.shadow.inherit = False

    title = spec.get("title")
    if title:
        box = textbox(slide, theme.MARGIN_L, theme.TITLE_TOP, theme.CONTENT_W, theme.TITLE_H)
        para = box.text_frame.paragraphs[0]
        rich_runs(para, title, size=theme.SIZE_TITLE, color=theme.BLUE_DEEP, bold=True)
        for run in para.runs:
            run.font.name = theme.FONT_HEAD

    kicker = spec.get("kicker")
    if kicker:
        box = textbox(slide, theme.MARGIN_L, theme.KICKER_TOP, theme.CONTENT_W, theme.KICKER_H)
        para = box.text_frame.paragraphs[0]
        rich_runs(para, kicker, size=theme.SIZE_KICKER, color=style.color, italic=True)

    rule = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        theme.MARGIN_L,
        theme.FOOTER_TOP - Inches(0.08),
        theme.SLIDE_WIDTH - theme.MARGIN_R,
        theme.FOOTER_TOP - Inches(0.08),
    )
    rule.line.color.rgb = theme.GREY_RULE
    rule.line.width = Pt(0.75)

    footer_text = spec.get("footer", doc.get("footer", ""))
    box = textbox(slide, theme.MARGIN_L, theme.FOOTER_TOP, theme.CONTENT_W, theme.FOOTER_H)
    para = box.text_frame.paragraphs[0]
    left = footer_text or style.label
    run = para.add_run()
    run.text = left
    style_run(run, size=theme.SIZE_FOOTER, color=theme.GREY_MUTED)

    number = textbox(
        slide,
        theme.SLIDE_WIDTH - theme.MARGIN_R - Inches(0.9),
        theme.FOOTER_TOP,
        Inches(0.9),
        theme.FOOTER_H,
    )
    para = number.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = str(deck.slide_count)
    style_run(run, size=theme.SIZE_FOOTER, color=theme.GREY_MUTED)


def body_top(spec: dict[str, Any]) -> Emu:
    """
    Where the content area starts, allowing for the optional kicker line.
    """
    return theme.BODY_TOP if spec.get("kicker") else theme.KICKER_TOP + Inches(0.12)


def render_bullet_list(frame, bullets: list[Any], *, base_size_map=None, first=True) -> None:
    """
    Fill a text frame with the deck's three-level bullet style.
    """
    sizes = base_size_map or theme.BULLET_SIZES
    for item in flatten_bullets(bullets):
        level, text, style_hint = item
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        para.space_after = Pt(8 if level == 0 else 4)
        para.space_before = Pt(6 if level == 0 else 0)
        para.line_spacing = 1.02
        size = sizes.get(level, theme.SIZE_BULLET_3)
        color = theme.GREY_TEXT if level else theme.BLUE_DEEP
        if style_hint == "muted":
            color = theme.GREY_MUTED
        mark = theme.BULLET_MARKS.get(level, "\u00b7")
        marker = para.add_run()
        marker.text = f"{mark}  "
        style_run(
            marker,
            size=Pt(size.pt - 2),
            color=theme.ACCENT if level == 0 else theme.BLUE_LIGHT,
            bold=True,
        )
        rich_runs(
            para,
            text,
            size=size,
            color=color,
            bold=(level == 0),
            italic=(style_hint == "muted"),
        )
        indent = theme.BULLET_INDENT.get(level, Inches(0.64)) + Inches(0.24)
        set_indent(para, indent, Inches(0.24))

def flatten_bullets(bullets: list[Any], level: int = 0) -> list[tuple[int, str, str | None]]:
    """
    Normalise the nested/flat bullet forms into ``(level, text, style)`` triples.
    """
    out: list[tuple[int, str, str | None]] = []
    for item in bullets or []:
        if isinstance(item, str):
            out.append((level, item, None))
            continue
        text = item.get("text", "")
        item_level = item.get("level", level)
        out.append((item_level, text, item.get("style")))
        for child in item.get("children", []) or []:
            if isinstance(child, str):
                out.append((item_level + 1, child, item.get("child_style")))
            else:
                out.extend(flatten_bullets([child], item_level + 1))
    return out


def layout_title(deck: Deck, slide, spec, doc) -> None:
    """
    The deck's opening slide: full-bleed panel, title, subtitle, byline.
    """
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), theme.SLIDE_WIDTH, theme.SLIDE_HEIGHT
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.BLUE_DEEP
    panel.line.fill.background()
    panel.shadow.inherit = False

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, theme.MARGIN_L, Inches(2.30), Inches(1.6), Inches(0.09)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False

    box = textbox(slide, theme.MARGIN_L, Inches(2.62), theme.CONTENT_W, Inches(1.7))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = spec.get("title", "")
    style_run(run, size=Pt(46), color=theme.WHITE, font=theme.FONT_HEAD, bold=True)

    subtitle = spec.get("subtitle")
    if subtitle:
        box = textbox(slide, theme.MARGIN_L, Inches(4.22), theme.CONTENT_W, Inches(1.0))
        para = box.text_frame.paragraphs[0]
        rich_runs(para, subtitle, size=Pt(19), color=theme.BLUE_LIGHT)

    for offset, line in enumerate(spec.get("byline", []) or []):
        box = textbox(
            slide, theme.MARGIN_L, Inches(5.55) + Inches(0.32) * offset, theme.CONTENT_W, Inches(0.32)
        )
        para = box.text_frame.paragraphs[0]
        rich_runs(para, line, size=Pt(12), color=theme.GREY_RULE)

    box = textbox(slide, theme.MARGIN_L, Inches(6.85), theme.CONTENT_W, Inches(0.4))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = spec.get("disclaimer", theme.DISCLAIMER)
    style_run(run, size=Pt(10), color=theme.GREY_MUTED, italic=True)


def layout_section(deck: Deck, slide, spec, doc) -> None:
    """
    A divider announcing a tree, with the list of what it contains.
    """
    style = theme.section_style(doc.get("section", "front"))
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), theme.SLIDE_WIDTH, theme.SLIDE_HEIGHT
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = style.color
    panel.line.fill.background()
    panel.shadow.inherit = False

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, theme.MARGIN_L, Inches(2.05), Inches(1.2), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False

    box = textbox(slide, theme.MARGIN_L, Inches(2.35), theme.CONTENT_W, Inches(1.1))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = spec.get("title", "")
    style_run(run, size=Pt(40), color=theme.WHITE, font=theme.FONT_HEAD, bold=True)

    if spec.get("subtitle"):
        box = textbox(slide, theme.MARGIN_L, Inches(3.52), Inches(9.2), Inches(0.9))
        para = box.text_frame.paragraphs[0]
        rich_runs(para, spec["subtitle"], size=Pt(17), color=theme.WHITE)

    items = spec.get("contents", []) or []
    for offset, item in enumerate(items):
        box = textbox(
            slide,
            theme.MARGIN_L,
            Inches(4.55) + Inches(0.36) * offset,
            theme.CONTENT_W,
            Inches(0.36),
        )
        para = box.text_frame.paragraphs[0]
        marker = para.add_run()
        marker.text = "\u25a0  "
        style_run(marker, size=Pt(13), color=theme.ACCENT, bold=True)
        rich_runs(para, item, size=Pt(14), color=theme.WHITE)


def layout_bullets(deck: Deck, slide, spec, doc) -> None:
    """
    The workhorse layout: a title, an optional kicker, and up to three bullet levels.
    """
    chrome(deck, slide, spec, doc)
    top = body_top(spec)
    box = textbox(slide, theme.MARGIN_L, top, theme.CONTENT_W, theme.FOOTER_TOP - top - Inches(0.2))
    render_bullet_list(box.text_frame, spec.get("bullets", []))

    if spec.get("takeaway"):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            theme.MARGIN_L,
            theme.FOOTER_TOP - Inches(0.92),
            theme.CONTENT_W,
            Inches(0.62),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = theme.BLUE_PALE
        panel.line.color.rgb = theme.BLUE_LIGHT
        panel.line.width = Pt(0.75)
        panel.shadow.inherit = False
        frame = panel.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.16)
        frame.margin_right = Inches(0.16)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        rich_runs(para, spec["takeaway"], size=Pt(13), color=theme.BLUE_DEEP)


def column_block(slide, left, top, width, height, block, accent) -> None:
    """
    Render one column of a two-column slide: heading, bullets, optional panel.
    """
    heading = block.get("heading")
    cursor = Emu(int(top))
    left = Emu(int(left))
    width = Emu(int(width))
    if heading:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, cursor, Inches(0.06), Inches(0.30))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False
        box = textbox(slide, left + Inches(0.16), cursor, width - Inches(0.16), Inches(0.34))
        para = box.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = heading
        style_run(run, size=Pt(15), color=theme.BLUE_DEEP, font=theme.FONT_HEAD, bold=True)
        cursor += Inches(0.46)
    box = textbox(slide, left, cursor, width, Emu(int(height - (cursor - top))))
    sizes = {0: Pt(14), 1: Pt(12), 2: Pt(11)}
    render_bullet_list(box.text_frame, block.get("bullets", []), base_size_map=sizes)


def layout_two_column(deck: Deck, slide, spec, doc) -> None:
    """
    Two independent bullet columns, for contrasts and before/after comparisons.
    """
    chrome(deck, slide, spec, doc)
    style = theme.section_style(doc.get("section", "front"))
    top = body_top(spec)
    height = theme.FOOTER_TOP - top - Inches(0.2)
    gutter = Inches(0.45)
    width = Emu(int((theme.CONTENT_W - gutter) / 2))
    column_block(slide, theme.MARGIN_L, top, width, height, spec.get("left", {}), style.color)
    column_block(
        slide,
        Emu(int(theme.MARGIN_L + width + gutter)),
        top,
        width,
        height,
        spec.get("right", {}),
        theme.ACCENT_DARK,
    )


def layout_table(deck: Deck, slide, spec, doc) -> None:
    """
    A styled table; column widths are given as fractions of the content width.
    """
    chrome(deck, slide, spec, doc)
    style = theme.section_style(doc.get("section", "front"))
    spec_table = spec.get("table", {})
    columns = spec_table.get("columns", [])
    rows = spec_table.get("rows", [])
    if not columns or not rows:
        deck.problem(doc["_path"], "table layout needs both columns and rows")
        return

    top = body_top(spec)
    available = theme.FOOTER_TOP - top - Inches(0.2)
    row_h = min(Inches(0.42), Emu(int(available / (len(rows) + 1))))
    height = Emu(int(row_h * (len(rows) + 1)))
    shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), theme.MARGIN_L, top, theme.CONTENT_W, height
    )
    table = shape.table
    table.first_row = True
    table.horz_banding = True

    fractions = spec_table.get("widths") or [1 / len(columns)] * len(columns)
    total = sum(fractions)
    for index, fraction in enumerate(fractions):
        table.columns[index].width = Emu(int(theme.CONTENT_W * fraction / total))

    size = Pt(spec_table.get("font_size", theme.SIZE_TABLE.pt))
    for index, header in enumerate(columns):
        cell = table.cell(0, index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = style.color
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(header)
        style_run(run, size=size, color=theme.WHITE, font=theme.FONT_HEAD, bold=True)

    for r_index, row in enumerate(rows, start=1):
        for c_index in range(len(columns)):
            value = str(row[c_index]) if c_index < len(row) else ""
            cell = table.cell(r_index, c_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.WHITE if r_index % 2 else theme.GREY_PANEL
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            rich_runs(para, value, size=size, color=theme.GREY_TEXT)


def layout_code(deck: Deck, slide, spec, doc) -> None:
    """
    A monospaced panel, for NodeSet fragments, JSON, C# and command lines.
    """
    chrome(deck, slide, spec, doc)
    top = body_top(spec)
    intro = spec.get("intro")
    if intro:
        box = textbox(slide, theme.MARGIN_L, top, theme.CONTENT_W, Inches(0.42))
        rich_runs(box.text_frame.paragraphs[0], intro, size=Pt(13), color=theme.GREY_TEXT)
        top = top + Inches(0.52)

    height = theme.FOOTER_TOP - top - Inches(0.2)
    if spec.get("outro"):
        height -= Inches(0.5)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, theme.MARGIN_L, top, theme.CONTENT_W, height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.BLUE_DEEP
    panel.line.fill.background()
    panel.shadow.inherit = False

    frame = panel.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.22)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.16)
    frame.margin_bottom = Inches(0.16)
    size = Pt(spec.get("font_size", theme.SIZE_CODE.pt))
    lines = (spec.get("code", "") or "").rstrip("\n").split("\n")
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.space_after = Pt(0)
        para.line_spacing = 1.06
        run = para.add_run()
        run.text = line or " "
        comment = line.lstrip().startswith(("#", "//", "<!--"))
        style_run(
            run,
            size=size,
            color=theme.BLUE_LIGHT if comment else theme.WHITE,
            font=theme.FONT_MONO,
            italic=comment,
        )

    if spec.get("outro"):
        box = textbox(
            slide, theme.MARGIN_L, top + height + Inches(0.12), theme.CONTENT_W, Inches(0.4)
        )
        rich_runs(box.text_frame.paragraphs[0], spec["outro"], size=Pt(12), color=theme.GREY_TEXT)


BOX_STYLES = {
    "primary": (theme.BLUE, theme.WHITE),
    "accent": (theme.ACCENT, theme.BLUE_DEEP),
    "pale": (theme.BLUE_PALE, theme.BLUE_DEEP),
    "muted": (theme.GREY_PANEL, theme.GREY_TEXT),
    "deep": (theme.BLUE_DEEP, theme.WHITE),
    "green": (theme.GREEN, theme.WHITE),
    "purple": (theme.SECTIONS["metaverse"].color, theme.WHITE),
}


def layout_diagram(deck: Deck, slide, spec, doc) -> None:
    """
    Boxes and connectors on a 12x6 grid, for architecture and flow pictures.

    Node geometry is given in grid units so the YAML stays readable; ``cols`` and
    ``rows`` set the grid resolution.
    """
    chrome(deck, slide, spec, doc)
    top = body_top(spec)
    height = theme.FOOTER_TOP - top - Inches(0.2)
    cols = spec.get("cols", 12)
    rows = spec.get("rows", 6)
    cell_w = theme.CONTENT_W / cols
    cell_h = height / rows

    shapes: dict[str, Any] = {}
    for node in spec.get("nodes", []) or []:
        fill, text_color = BOX_STYLES.get(node.get("style", "pale"), BOX_STYLES["pale"])
        left = Emu(int(theme.MARGIN_L + cell_w * node.get("x", 0)))
        node_top = Emu(int(top + cell_h * node.get("y", 0)))
        width = Emu(int(cell_w * node.get("w", 3)))
        node_h = Emu(int(cell_h * node.get("h", 1)))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if node.get("round", True) else MSO_SHAPE.RECTANGLE,
            left,
            node_top,
            width,
            node_h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = theme.BLUE_LIGHT if node.get("style") == "pale" else fill
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        frame.margin_top = Inches(0.04)
        frame.margin_bottom = Inches(0.04)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        rich_runs(
            para,
            node.get("text", ""),
            size=Pt(node.get("font_size", 12)),
            color=text_color,
            bold=node.get("bold", True),
        )
        for extra in node.get("lines", []) or []:
            para = frame.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            para.space_before = Pt(1)
            rich_runs(para, extra, size=Pt(node.get("line_size", 10)), color=text_color)
        if node.get("id"):
            shapes[node["id"]] = shape

    for arrow in spec.get("arrows", []) or []:
        source = shapes.get(arrow.get("from"))
        target = shapes.get(arrow.get("to"))
        if source is None or target is None:
            deck.problem(doc["_path"], f"arrow references unknown node: {arrow}")
            continue
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(source.left + source.width // 2),
            Emu(source.top + source.height // 2),
            Emu(target.left + target.width // 2),
            Emu(target.top + target.height // 2),
        )
        connector.line.color.rgb = theme.GREY_MUTED
        connector.line.width = Pt(arrow.get("width", 1.25))
        connector.begin_connect(source, arrow.get("from_site", 3))
        connector.end_connect(target, arrow.get("to_site", 1))
        if arrow.get("label"):
            mid_x = (source.left + source.width // 2 + target.left + target.width // 2) // 2
            mid_y = (source.top + source.height // 2 + target.top + target.height // 2) // 2
            box = textbox(
                slide, Emu(int(mid_x - Inches(0.85))), Emu(int(mid_y - Inches(0.16))), Inches(1.7), Inches(0.3)
            )
            para = box.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = arrow["label"]
            style_run(run, size=Pt(9), color=theme.GREY_MUTED, italic=True)

    for caption in spec.get("captions", []) or []:
        box = textbox(
            slide,
            Emu(int(theme.MARGIN_L + cell_w * caption.get("x", 0))),
            Emu(int(top + cell_h * caption.get("y", 0))),
            Emu(int(cell_w * caption.get("w", 4))),
            Inches(0.4),
        )
        para = box.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if caption.get("center") else PP_ALIGN.LEFT
        rich_runs(
            para,
            caption.get("text", ""),
            size=Pt(caption.get("font_size", 11)),
            color=theme.GREY_MUTED,
            italic=True,
        )


def layout_statement(deck: Deck, slide, spec, doc) -> None:
    """
    One sentence, large, for the claim a section turns on.
    """
    style = theme.section_style(doc.get("section", "front"))
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), theme.SLIDE_WIDTH, theme.SLIDE_HEIGHT
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.GREY_PANEL
    panel.line.fill.background()
    panel.shadow.inherit = False

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(0.16), theme.SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.color
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = textbox(slide, Inches(1.25), Inches(2.35), Inches(10.9), Inches(2.4))
    frame = box.text_frame
    para = frame.paragraphs[0]
    rich_runs(para, spec.get("statement", ""), size=Pt(30), color=theme.BLUE_DEEP, bold=True)
    for run in para.runs:
        run.font.name = theme.FONT_HEAD
    if spec.get("attribution"):
        para = frame.add_paragraph()
        para.space_before = Pt(18)
        rich_runs(para, spec["attribution"], size=Pt(14), color=theme.GREY_MUTED, italic=True)


def layout_demo(deck: Deck, slide, spec, doc) -> None:
    """
    The demo slide: what you will see, the moving parts, the command, what it proves.
    """
    chrome(deck, slide, spec, doc)
    top = body_top(spec)

    state = spec.get("state", "master")
    badge_text = {
        "master": "RUNS ON master",
        "branch": "RUNS ON A BRANCH",
        "walkthrough": "WALKTHROUGH ONLY",
    }.get(state, state.upper())
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        theme.SLIDE_WIDTH - theme.MARGIN_R - Inches(2.05),
        theme.TITLE_TOP + Inches(0.06),
        Inches(2.05),
        Inches(0.34),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = theme.DEMO_STATE_COLORS.get(state, theme.GREY_MUTED)
    badge.line.fill.background()
    badge.shadow.inherit = False
    frame = badge.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = badge_text
    style_run(run, size=Pt(9.5), color=theme.WHITE, font=theme.FONT_HEAD, bold=True)

    gutter = Inches(0.4)
    col_w = Emu(int((theme.CONTENT_W - gutter) / 2))
    body_h = Inches(2.55)
    column_block(
        slide,
        theme.MARGIN_L,
        top,
        col_w,
        body_h,
        {"heading": "What you will see", "bullets": spec.get("see", [])},
        theme.GREEN,
    )
    column_block(
        slide,
        Emu(int(theme.MARGIN_L + col_w + gutter)),
        top,
        col_w,
        body_h,
        {"heading": "The moving parts", "bullets": spec.get("parts", [])},
        theme.BLUE,
    )

    run_top = Emu(int(top + body_h + Inches(0.12)))
    commands = spec.get("run", []) or []
    if commands:
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            theme.MARGIN_L,
            run_top,
            theme.CONTENT_W,
            Emu(int(Inches(0.42) + Inches(0.26) * (len(commands) - 1))),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = theme.BLUE_DEEP
        panel.line.fill.background()
        panel.shadow.inherit = False
        frame = panel.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.18)
        frame.margin_top = Inches(0.08)
        frame.margin_bottom = Inches(0.08)
        for index, command in enumerate(commands):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.space_after = Pt(0)
            run = para.add_run()
            # The panel is already monospaced, so a wrapping backtick is never wanted.
            run.text = str(command).strip().strip("`")
            style_run(run, size=Pt(11.5), color=theme.WHITE, font=theme.FONT_MONO)
        run_top = Emu(int(run_top + panel.height + Inches(0.12)))

    if spec.get("proves"):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            theme.MARGIN_L,
            run_top,
            theme.CONTENT_W,
            Emu(int(theme.FOOTER_TOP - run_top - Inches(0.18))),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = theme.BLUE_PALE
        panel.line.color.rgb = theme.BLUE_LIGHT
        panel.line.width = Pt(0.75)
        panel.shadow.inherit = False
        frame = panel.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.16)
        frame.margin_right = Inches(0.16)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        label = para.add_run()
        label.text = "What it proves:  "
        style_run(label, size=Pt(12.5), color=theme.BLUE, font=theme.FONT_HEAD, bold=True)
        rich_runs(para, spec["proves"], size=Pt(12.5), color=theme.BLUE_DEEP)


RENDERERS = {
    "title": layout_title,
    "section": layout_section,
    "bullets": layout_bullets,
    "two-column": layout_two_column,
    "table": layout_table,
    "code": layout_code,
    "diagram": layout_diagram,
    "statement": layout_statement,
    "demo": layout_demo,
}


def build(docs: list[dict[str, Any]]) -> Deck:
    """
    Render every slide in every content document into one presentation.
    """
    deck = Deck(prs=new_presentation())
    for doc in docs:
        path = doc.get("_path", "?")
        if "section" not in doc:
            deck.problem(path, "missing 'section'")
        for index, spec in enumerate(doc.get("slides", []) or []):
            layout = spec.get("layout", "bullets")
            if layout not in LAYOUTS:
                deck.problem(path, f"slide {index}: unknown layout '{layout}'")
                continue
            slide = add_slide(deck)
            RENDERERS[layout](deck, slide, spec, doc)
            set_notes(slide, spec.get("notes"))
            if layout not in {"title", "section", "statement"} and not spec.get("notes"):
                deck.problem(path, f"slide {index} ('{spec.get('title', layout)}') has no notes")
    return deck


def print_outline(docs: list[dict[str, Any]]) -> None:
    """
    Print the deck's running order: one line per slide, grouped by content file.
    """
    number = 0
    current_section = None
    for doc in docs:
        section = doc.get("section", "?")
        if section != current_section:
            current_section = section
            print(f"\n=== {theme.section_style(section).label} ===")
        print(f"  {doc['_path']}  (order {doc.get('order', '?')})")
        for spec in doc.get("slides", []) or []:
            number += 1
            layout = spec.get("layout", "bullets")
            title = spec.get("title") or spec.get("statement") or ""
            title = " ".join(str(title).split())[:78]
            print(f"    {number:>3}  {layout:<11} {title}")
    print(f"\n{number} slides")


def main(argv: list[str] | None = None) -> int:
    """
    Command-line entry point.
    """
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--check", action="store_true", help="validate only, write nothing")
    parser.add_argument("--outline", action="store_true", help="print the running order and exit")
    parser.add_argument("--strict", action="store_true", help="treat content problems as failure")
    args = parser.parse_args(argv)

    docs, errors = load_content()
    if not docs and not errors:
        print(f"no content found under {CONTENT_DIR}", file=sys.stderr)
        return 1

    if args.outline:
        for problem in errors:
            print(f"warning: {problem}", file=sys.stderr)
        print_outline(docs)
        return 1 if (args.strict and errors) else 0

    deck = build(docs)
    deck.problems = errors + deck.problems
    for problem in deck.problems:
        print(f"warning: {problem}", file=sys.stderr)

    if not args.check:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        deck.prs.save(args.out)
        print(f"wrote {args.out} \u2014 {deck.slide_count} slides from {len(docs)} content files")
    else:
        print(f"checked {deck.slide_count} slides from {len(docs)} content files")

    return 1 if (args.strict and deck.problems) else 0


if __name__ == "__main__":
    raise SystemExit(main())
