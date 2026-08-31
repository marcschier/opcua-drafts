"""Mermaid flowchart -> an editable PowerPoint object plus a raster preview.

OPC 20020 Guideline 1 requires figures to be embedded document objects (PowerPoint,
Excel or Visio) and forbids inline Word drawing objects, so a rendered bitmap alone
would not conform. The diagram is therefore rebuilt as real PowerPoint shapes an editor
can move, and the bitmap exists only as the OLE preview Word displays.

Only the flowchart subset the drafts use is understood; an unrecognised construct raises
rather than being dropped.
"""

import html
import re

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

EMU_PER_PX = 9525

# The OPC UA graphical notation draws a type with a drop shadow and an abstract type in
# italics, so the renderer has to know which NodeClass a box stands for. Mermaid carries
# that as `:::class`, which is standard syntax and renders on GitHub through `classDef`.
TYPE_CLASSES = {'objecttype', 'variabletype', 'referencetype', 'datatype', 'eventtype'}
INSTANCE_CLASSES = {'object', 'variable', 'view', 'method'}
NODE_CLASSES = TYPE_CLASSES | INSTANCE_CLASSES | {'abstract', 'placeholder'}

CLASS_SUFFIX = r'(?::::(?P<cls>[A-Za-z0-9_]+(?:,[A-Za-z0-9_]+)*))?'
NODE_RE = re.compile(r'^(?P<id>[A-Za-z0-9_]+)\s*(?P<open>\[\[|\[|\(\(|\(|\{)'
                     r'(?P<label>.*?)'
                     r'(?P<close>\]\]|\]|\)\)|\)|\})' + CLASS_SUFFIX + r'\s*$')
# A node may also be referenced by bare id with a class attached, e.g. `A:::object --> B`.
BARE_ID_RE = re.compile(r'^(?P<id>[A-Za-z0-9_]+)' + CLASS_SUFFIX + r'\s*$')
# Mermaid edge tokens. The labelled forms must precede the plain ones: `-->` would
# otherwise match the head of `-->|label|` and leave `|label| B` looking like a node.
# The `A -- text --> B` form needs its own alternative for the same reason — without it
# the split happens at the trailing `-->` and `A -- text` is taken for a node id.
EDGE_SPLIT_RE = re.compile(
    r'(-{2,}>\|[^|]*\|'
    r'|={2,}>\|[^|]*\|'
    r'|-\.+->\|[^|]*\|'
    r'|-\.[^.|]*\.->'
    r'|-\.->'
    r'|-{2,}\s[^|]*?\s-{2,}>'
    r'|={2,}\s[^|]*?\s={2,}>'
    r'|-{2,}>'
    r'|={2,}>)')
EDGE_ID_RE = re.compile(r'^[A-Za-z0-9_]+$')
EDGE_MID_LABEL_RE = re.compile(r'^[-=]{2,}\s+(?P<label>.*?)\s+[-=]{2,}>$')
SUBGRAPH_RE = re.compile(r'^subgraph\s+(?P<id>[A-Za-z0-9_]+)\s*(?:\["?(?P<label>.*?)"?\])?\s*$')


class Node:
    def __init__(self, nid, label, shape='box', cls=()):
        self.id = nid
        self.label = label
        self.shape = shape
        # Mermaid `:::name` classes. They carry the OPC UA NodeClass of the Node the box
        # stands for, which decides whether it is drawn as a plain rectangle, a rounded
        # one, or a type box with a drop shadow.
        self.cls = tuple(cls)
        self.subgraph = None
        self.layer = 0
        self.x = 0
        self.y = 0
        self.w = 0
        self.h = 0

    @property
    def is_type(self):
        return bool(TYPE_CLASSES & set(self.cls))

    @property
    def is_abstract(self):
        return 'abstract' in self.cls


class Edge:
    def __init__(self, src, dst, label=None, dashed=False, thick=False):
        self.src = src
        self.dst = dst
        self.label = label
        self.dashed = dashed
        # `==>` is HasTypeDefinition, which OPC UA draws with a solid head. Without
        # recording it the thick form parsed identically to `-->` and the two references
        # were indistinguishable in the rendering.
        self.thick = thick
        # Filled in by `_place_edge_labels` once the nodes have positions.
        self.lx = self.ly = self.lw = 0


class Graph:
    def __init__(self):
        self.nodes = {}
        self.edges = []
        self.subgraphs = {}
        # Mermaid nests subgraphs; a node records only its innermost one, so the nesting
        # has to be kept separately or an outer frame is drawn beside its own child.
        self.subgraph_parent = {}
        self.order = []

    def node(self, nid, label=None, shape='box', cls=()):
        n = self.nodes.get(nid)
        if n is None:
            n = Node(nid, label if label is not None else nid, shape, cls)
            self.nodes[nid] = n
            self.order.append(nid)
        else:
            if label is not None and n.label == n.id:
                n.label = label
            if cls and not n.cls:
                n.cls = tuple(cls)
        return n

    @property
    def opcua(self):
        """Whether this diagram uses the OPC UA notation.

        The notation is opt-in, decided by any node carrying a `:::` NodeClass. Without
        this every existing architecture and sequence figure would silently change
        appearance the moment the renderer learned the notation, and that churn would be
        indistinguishable from an intended edit in review.
        """
        return any(n.cls for n in self.nodes.values())


def _clean(label):
    label = label.strip().strip('"')
    label = label.replace('<br/>', '\n').replace('<br>', '\n')
    # Mermaid needs `<` and `>` escaped, but a placeholder BrowseName such as
    # `<WoTAssetName>` has to reach the figure as the model spells it, not as `&lt;...`.
    return html.unescape(label)


def parse(source):
    """Parse a Mermaid diagram; dispatches on the diagram keyword."""
    head = next((ln.strip() for ln in source.splitlines() if ln.strip()), '')
    if head.startswith('sequenceDiagram'):
        return parse_sequence(source)
    if head.startswith('classDiagram'):
        return parse_class_diagram(source)
    if head.startswith('stateDiagram'):
        return parse_state_diagram(source)
    return parse_flowchart(source)


STATE_EDGE_RE = re.compile(
    r'^(?P<a>\[\*\]|[A-Za-z0-9_]+)\s*-->\s*(?P<b>\[\*\]|[A-Za-z0-9_]+)'
    r'\s*(?::\s*(?P<label>.*))?$')
STATE_ALIAS_RE = re.compile(r'^state\s+"(?P<label>[^"]*)"\s+as\s+(?P<id>[A-Za-z0-9_]+)$')

# Mermaid's start and end pseudostates are both spelled `[*]`; which one it means
# depends on the side of the arrow it sits on.
STATE_START = '__start'
STATE_END = '__end'


def parse_state_diagram(source):
    """A state diagram, laid out with the flowchart machinery.

    States are nodes and transitions are labelled edges, so once `[*]` is resolved into
    an explicit start or end node the existing layered layout renders it unchanged.
    """
    g = Graph()
    for raw in source.splitlines():
        line = raw.strip()
        if not line or line.startswith('%%'):
            continue
        if line.startswith(('stateDiagram', 'direction')):
            continue
        alias = STATE_ALIAS_RE.match(line)
        if alias:
            g.node(alias.group('id'), _clean(alias.group('label')), 'round')
            continue
        m = STATE_EDGE_RE.match(line)
        if not m:
            raise ValueError('unsupported mermaid state line: %r' % raw)
        src = _state_node(g, m.group('a'), STATE_START)
        dst = _state_node(g, m.group('b'), STATE_END)
        label = m.group('label')
        g.edges.append(Edge(src, dst, _clean(label) if label else None))
    return g


def _state_node(g, token, pseudo_id):
    if token == '[*]':
        return g.node(pseudo_id, 'start' if pseudo_id == STATE_START else 'end',
                      'circle').id
    return g.node(token, token, 'round').id


CLASS_OPEN_RE = re.compile(r'^class\s+(?P<id>[A-Za-z0-9_]+)\s*\{$')
CLASS_ONELINE_RE = re.compile(r'^class\s+(?P<id>[A-Za-z0-9_]+)\s*$')
CLASS_REL_RE = re.compile(
    r'^(?P<a>[A-Za-z0-9_]+)\s*'
    r'(?P<rel><\|--|--\|>|\.\.\|>|<\|\.\.|\*--|o--|-->|<--|\.\.>|\.\.|--)\s*'
    r'(?P<b>[A-Za-z0-9_]+)\s*(?::\s*(?P<label>.*))?$')


def parse_class_diagram(source):
    """A UML class diagram, laid out with the same machinery as a flowchart.

    Each class becomes a node whose label is its name over its members, and each
    relation becomes an edge; inheritance is drawn from the subtype to its base so the
    layering puts bases above their subtypes.
    """
    g = Graph()
    members = {}
    current = None
    for raw in source.splitlines():
        line = raw.strip()
        if not line or line.startswith('%%'):
            continue
        if line.startswith('classDiagram') or line.startswith('direction'):
            continue
        if line == '}':
            current = None
            continue
        m = CLASS_OPEN_RE.match(line)
        if m:
            current = m.group('id')
            g.node(current)
            members.setdefault(current, [])
            continue
        if current:
            members[current].append(line.lstrip('+-#').strip())
            continue
        m = CLASS_ONELINE_RE.match(line)
        if m:
            g.node(m.group('id'))
            continue
        m = CLASS_REL_RE.match(line)
        if m:
            a, b = g.node(m.group('a')).id, g.node(m.group('b')).id
            rel = m.group('rel')
            # `B --|> A` and `A <|-- B` both mean "B is a subtype of A".
            if rel in ('--|>', '*--', 'o--'):
                src, dst = a, b
            elif rel == '<|--':
                src, dst = b, a
            else:
                src, dst = a, b
            g.edges.append(Edge(src, dst, _clean(m.group('label') or '') or None,
                                rel in ('..>',)))
            continue
        raise ValueError('unsupported mermaid class-diagram line: %r' % raw)

    for nid, lines in members.items():
        if lines:
            g.nodes[nid].label = nid + '\n' + '\n'.join(lines)
    return g


def parse_flowchart(source):
    """Parse the `flowchart` subset used by the drafts."""
    g = Graph()
    stack = []
    for raw in source.splitlines():
        line = raw.strip()
        if not line or line.startswith('%%'):
            continue
        if line.startswith(('flowchart', 'graph', 'direction')):
            continue
        # `classDef` and `class` style the diagram for GitHub's renderer. The shapes here
        # are driven by the `:::` classes on the nodes themselves, so the style lines are
        # carried by the markdown but are not part of the model this renderer draws.
        if line.startswith(('classDef', 'class ', 'linkStyle', 'style ')):
            continue
        if line == 'end':
            if stack:
                stack.pop()
            continue
        sub = SUBGRAPH_RE.match(line)
        if sub:
            sid = sub.group('id')
            g.subgraphs[sid] = _clean(sub.group('label') or sid)
            g.subgraph_parent[sid] = stack[-1] if stack else None
            stack.append(sid)
            continue
        if EDGE_SPLIT_RE.search(line):
            _parse_edge_chain(g, line, stack[-1] if stack else None)
            continue
        m = NODE_RE.match(line)
        if m:
            n = g.node(m.group('id'), _clean(m.group('label')),
                       _shape_for(m.group('open')), _classes_of(m))
            if stack and n.subgraph is None:
                n.subgraph = stack[-1]
            continue
        raise ValueError('unsupported mermaid line: %r' % raw)
    _resolve_cluster_endpoints(g)
    return g


def _resolve_cluster_endpoints(g):
    """An edge that names a subgraph attaches to that subgraph, not to a phantom node.

    Mermaid lets an edge end on a cluster id and draws it to the cluster border. Creating
    a node for it instead put an extra box in the figure captioned with the cluster's own
    id — a box that appears nowhere in the source. The edge is re-pointed at the first
    member of the cluster, which is the node the cluster's border stands for.
    """
    representative = {}
    for nid in g.order:
        sid = g.nodes[nid].subgraph
        while sid:
            representative.setdefault(sid, nid)
            sid = g.subgraph_parent.get(sid)
    phantom = {nid for nid in list(g.nodes)
               if nid in g.subgraphs and g.nodes[nid].label == nid
               and nid in representative}
    if not phantom:
        return
    for e in g.edges:
        e.src = representative.get(e.src, e.src) if e.src in phantom else e.src
        e.dst = representative.get(e.dst, e.dst) if e.dst in phantom else e.dst
    g.edges[:] = [e for e in g.edges if e.src != e.dst]
    for nid in phantom:
        del g.nodes[nid]
        g.order.remove(nid)


def _shape_for(opener):
    return {'[': 'box', '[[': 'subroutine', '(': 'round', '((': 'circle',
            '{': 'diamond'}.get(opener, 'box')


def _classes_of(match):
    raw = match.groupdict().get('cls')
    if not raw:
        return ()
    # Mermaid allows one `:::` per node, with several classes separated by commas.
    # Chaining `:::a:::b` is a parse error in Mermaid itself, so it cannot be accepted
    # here either — the figure has to render on GitHub as well as in the Word build.
    names = tuple(p for p in raw.split(',') if p)
    unknown = [p for p in names if p.lower() not in NODE_CLASSES]
    if unknown:
        # An unknown class is almost always a typo in a NodeClass name, and ignoring it
        # would draw an instance box where a type box belongs.
        raise ValueError('unknown mermaid node class: %s' % ', '.join(unknown))
    return tuple(p.lower() for p in names)


def _parse_edge_chain(g, line, subgraph):
    parts = EDGE_SPLIT_RE.split(line)
    endpoints = parts[0::2]
    connectors = parts[1::2]
    ids = []
    for token in endpoints:
        token = token.strip()
        m = NODE_RE.match(token)
        if m:
            n = g.node(m.group('id'), _clean(m.group('label')),
                       _shape_for(m.group('open')), _classes_of(m))
        else:
            # Without this an edge form the splitter does not know silently becomes a
            # node whose id is the unparsed text, and the diagram renders as nonsense
            # instead of failing the build.
            bare = BARE_ID_RE.match(token)
            if not bare:
                raise ValueError('unsupported mermaid edge endpoint: %r' % token)
            n = g.node(bare.group('id'), cls=_classes_of(bare))
        if subgraph and n.subgraph is None:
            n.subgraph = subgraph
        ids.append(n.id)
    for i, conn in enumerate(connectors):
        label = None
        dashed = conn.startswith('-.')
        thick = conn.startswith('=')
        lm = re.search(r'\|([^|]*)\|', conn)
        if lm:
            label = _clean(lm.group(1))
        dm = re.match(r'-\.(.+)\.->$', conn)
        if dm:
            label = _clean(dm.group(1))
        mid = EDGE_MID_LABEL_RE.match(conn)
        if mid:
            label = _clean(mid.group('label'))
        g.edges.append(Edge(ids[i], ids[i + 1], label, dashed, thick))


# --------------------------------------------------------------------------- layout

BOX_W = 230
BOX_H = 56
H_GAP = 40
V_GAP = 58
MARGIN = 26
SUBGRAPH_PAD = 18
TITLE_H = 24
# The widest row a figure may have. A page is about six boxes wide at a readable size.
MAX_ROW = 6


def _acyclic_edges(g):
    """The edges that may constrain layering: cycle-closing back edges removed.

    A state machine legitimately loops (Failed -> Loading), and longest-path layering
    over a cycle diverges — each pass pushes the nodes one layer further down until the
    guard trips, which produced a canvas hundreds of inches tall. Back edges are still
    drawn; they just do not decide which layer a node sits on.
    """
    adjacency = {}
    for e in g.edges:
        adjacency.setdefault(e.src, []).append(e.dst)
    state = {}
    back = set()

    def visit(start):
        stack = [(start, iter(adjacency.get(start, ())))]
        state[start] = 1
        while stack:
            nid, children = stack[-1]
            for dst in children:
                if state.get(dst) == 1:
                    back.add((nid, dst))
                elif state.get(dst) is None:
                    state[dst] = 1
                    stack.append((dst, iter(adjacency.get(dst, ()))))
                    break
            else:
                state[nid] = 2
                stack.pop()

    for nid in g.order:
        if state.get(nid) is None:
            visit(nid)
    return [e for e in g.edges if (e.src, e.dst) not in back]


def _cluster_bands(g, forward):
    """Layer assignment that keeps every subgraph a contiguous band, or None.

    Mermaid draws a subgraph as a frame around its members. Plain longest-path layering
    can interleave two subgraphs' members, and the frames then overlap — there is no
    honest way to draw that, and the old code papered over it by shoving nodes sideways
    until the canvas was metres wide. Ranking the clusters first, then the nodes inside
    each cluster, makes every frame a solid block.

    Returns None when the graph has no subgraph, so a cluster-free diagram keeps exactly
    the layering it had before.
    """
    if not any(n.subgraph for n in g.nodes.values()):
        return None

    # An ungrouped node is its own singleton cluster, so it stays free to sit wherever
    # its edges put it instead of being herded in with every other ungrouped node. A
    # nested subgraph shares its outermost ancestor's band, so the parent frame stays
    # contiguous and can be drawn around its child.
    cluster_of = {nid: (_root_cluster(g, n.subgraph) if n.subgraph else ('\x00' + nid))
                  for nid, n in g.nodes.items()}
    members = {}
    for nid in g.order:
        members.setdefault(cluster_of[nid], []).append(nid)

    inner = [e for e in forward if cluster_of[e.src] == cluster_of[e.dst]]
    outer = [(cluster_of[e.src], cluster_of[e.dst]) for e in forward
             if cluster_of[e.src] != cluster_of[e.dst]]

    rank = _longest_path({c: None for c in members}, outer)
    local = _longest_path({nid: None for nid in g.nodes},
                          [(e.src, e.dst) for e in inner])

    height = {}
    for c, ids in members.items():
        height[c] = max(local[i] for i in ids) + 1
    base = {}
    offset = 0
    for r in sorted({rank[c] for c in members}):
        at_rank = [c for c in members if rank[c] == r]
        for c in at_rank:
            base[c] = offset
        offset += max(height[c] for c in at_rank)
    return {nid: base[cluster_of[nid]] + local[nid] for nid in g.nodes}


def _longest_path(keys, edges):
    """Longest-path ranking over an acyclic edge list."""
    rank = {k: 0 for k in keys}
    changed = True
    guard = 0
    while changed and guard < 200:
        changed = False
        guard += 1
        for src, dst in edges:
            if rank[dst] < rank[src] + 1:
                rank[dst] = rank[src] + 1
                changed = True
    return rank


def layout(g):
    """Longest-path layering, then centre each layer horizontally."""
    forward = _acyclic_edges(g)
    layer = _cluster_bands(g, forward)
    if layer is None:
        layer = _longest_path({nid: None for nid in g.nodes},
                              [(e.src, e.dst) for e in forward])
    for nid, lv in layer.items():
        g.nodes[nid].layer = lv

    layers = {}
    for nid in g.order:
        layers.setdefault(g.nodes[nid].layer, []).append(nid)
    # Members of a subgraph go on the left of their layer, so the subgraph frame does
    # not end up drawn around nodes that do not belong to it.
    for lv in layers:
        layers[lv].sort(key=lambda nid: (g.nodes[nid].subgraph is None,))

    # A layer of a hub-and-spoke model can hold a dozen nodes, and one row of a dozen
    # boxes is four times wider than the page it has to print on. Wide layers wrap, so a
    # figure grows downwards — where a document has room — instead of sideways.
    rows = []
    for lv in sorted(layers):
        ids = layers[lv]
        for start in range(0, len(ids), MAX_ROW):
            rows.append(ids[start:start + MAX_ROW])

    width = max((len(r) for r in rows), default=1)
    total_w = width * BOX_W + (width - 1) * H_GAP
    y = MARGIN + TITLE_H
    for ids in rows:
        row_w = len(ids) * BOX_W + (len(ids) - 1) * H_GAP
        x0 = MARGIN + (total_w - row_w) / 2
        for i, nid in enumerate(ids):
            n = g.nodes[nid]
            n.x = x0 + i * (BOX_W + H_GAP)
            n.y = y
            n.w = BOX_W
            n.h = BOX_H
        y += BOX_H + V_GAP

    _push_outsiders_clear(g)
    _place_edge_labels(g)
    frames = subgraph_bounds(g).values()
    left = min([n.x for n in g.nodes.values()] + [b[0] for b in frames] or [MARGIN])
    top = min([n.y for n in g.nodes.values()] + [b[1] for b in frames] or [MARGIN])
    if left < MARGIN or top < MARGIN:
        # A frame drawn around an outer subgraph reaches above and left of its members;
        # without this shift the outermost border is cropped off the slide.
        dx, dy = max(0, MARGIN - left), max(0, MARGIN - top)
        for n in g.nodes.values():
            n.x += dx
            n.y += dy
        for e in g.edges:
            e.lx += dx
            e.ly += dy
        frames = subgraph_bounds(g).values()
    right = max([n.x + n.w for n in g.nodes.values()]
                + [e.lx + e.lw for e in g.edges if e.label]
                + [b[2] for b in frames] or [total_w])
    bottom = max([n.y + n.h for n in g.nodes.values()]
                 + [e.ly + LABEL_H for e in g.edges if e.label]
                 + [b[3] for b in frames] or [0])
    return int(right + MARGIN), int(bottom + MARGIN)


LABEL_H = 13
LABEL_FONT_PX = 8


def _place_edge_labels(g):
    """Give every edge label a position clear of other labels and of the node boxes.

    Labels used to be drawn at the raw midpoint of their edge. Two edges running between
    the same pair of rows share that midpoint almost exactly, so the text overprinted;
    and a midpoint that lands on a node box is painted over when the boxes are drawn on
    top. Both produce a figure needing manual repair, which is what this pipeline exists
    to avoid. Colliding labels are stacked down the gap between the rows instead.
    """
    font = _font(LABEL_FONT_PX)
    obstacles = [(n.x, n.y, n.x + n.w, n.y + n.h) for n in g.nodes.values()]
    for e in g.edges:
        if not e.label:
            continue
        a, b = g.nodes[e.src], g.nodes[e.dst]
        width = _text_width(font, e.label)
        x = (a.x + a.w / 2 + b.x + b.w / 2) / 2 + 4
        y = (a.y + a.h + b.y) / 2 - 6
        for _ in range(40):
            if not any(_boxes_overlap((x, y, x + width, y + LABEL_H), o) for o in obstacles):
                break
            y += LABEL_H + 2
        obstacles.append((x, y, x + width, y + LABEL_H))
        e.lx, e.ly, e.lw = x, y, width


def _text_width(font, text):
    try:
        return font.getbbox(text)[2]
    except AttributeError:
        return int(len(text) * LABEL_FONT_PX * 0.62)


def _boxes_overlap(a, b):
    return not (a[2] <= b[0] or a[0] >= b[2] or a[3] <= b[1] or a[1] >= b[3])


def _push_outsiders_clear(g):
    """Move nodes that belong to no part of a subgraph out of that subgraph's frame."""
    for _ in range(4):
        bounds = subgraph_bounds(g)
        moved = False
        for n in g.nodes.values():
            for sid, (x0, y0, x1, y1) in bounds.items():
                if _within(g, n.subgraph, sid):
                    continue
                if n.x + n.w <= x0 or n.x >= x1 or n.y + n.h <= y0 or n.y >= y1:
                    continue
                n.x = x1 + H_GAP
                moved = True
        if moved:
            _separate_rows(g)
        else:
            break


def _separate_rows(g):
    """Pull apart nodes that ended up on the same row at the same x.

    Every outsider is pushed clear of a frame to the same coordinate, so two nodes evicted
    from the same frame landed exactly on top of each other and one of them vanished from
    the figure.
    """
    rows = {}
    for n in g.nodes.values():
        rows.setdefault(round(n.y), []).append(n)
    for row in rows.values():
        row.sort(key=lambda n: (n.x, n.id))
        for prev, node in zip(row, row[1:]):
            minimum = prev.x + prev.w + H_GAP
            if node.x < minimum:
                node.x = minimum


def _within(g, node_subgraph, frame):
    """True when a node sits inside a frame, directly or through a nested subgraph."""
    seen = set()
    sid = node_subgraph
    while sid and sid not in seen:
        if sid == frame:
            return True
        seen.add(sid)
        sid = g.subgraph_parent.get(sid)
    return False


def subgraph_bounds(g):
    """The frame of every subgraph, covering the nodes of any subgraph nested inside it.

    A frame that only covered its direct members would be drawn next to its own child
    instead of around it. Outer frames get proportionally more padding so each nested
    frame sits visibly inside its parent.
    """
    depth = {sid: _nesting_depth(g, sid) for sid in g.subgraphs}
    deepest = max(depth.values(), default=0)
    bounds = {}
    for n in g.nodes.values():
        sid = n.subgraph
        while sid:
            b = bounds.setdefault(sid, [1e9, 1e9, -1e9, -1e9])
            b[0] = min(b[0], n.x)
            b[1] = min(b[1], n.y)
            b[2] = max(b[2], n.x + n.w)
            b[3] = max(b[3], n.y + n.h)
            sid = g.subgraph_parent.get(sid)
    for sid, b in bounds.items():
        rings = deepest - depth.get(sid, 0) + 1
        b[0] -= SUBGRAPH_PAD * rings
        b[1] -= SUBGRAPH_PAD * rings + TITLE_H * rings
        b[2] += SUBGRAPH_PAD * rings
        b[3] += SUBGRAPH_PAD * rings
    return bounds


def _nesting_depth(g, sid):
    depth = 0
    seen = set()
    while g.subgraph_parent.get(sid) and sid not in seen:
        seen.add(sid)
        sid = g.subgraph_parent[sid]
        depth += 1
    return depth


def _root_cluster(g, sid):
    seen = set()
    while g.subgraph_parent.get(sid) and sid not in seen:
        seen.add(sid)
        sid = g.subgraph_parent[sid]
    return sid


# --------------------------------------------------------------------------- pptx

BORDER = RGBColor(0x44, 0x44, 0x44)
FILL = RGBColor(0xEE, 0xF3, 0xFA)
SUBGRAPH_FILL = RGBColor(0xF7, 0xF7, 0xF7)
TEXT = RGBColor(0x11, 0x11, 0x11)
# The OPC UA notation draws a type as a box standing on a grey offset block.
TYPE_SHADOW = 'BFBFBF'

_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Measured from the template's own figure (`Microsoft_PowerPoint_Slide.sldx` inside
# OPC 20020): every text run is Arial, the boxes are plain `rect`, and the connectors are
# a mix of `straightConnector1` and `bentConnector3`. Matching the typeface and the
# orthogonal routing is what makes a figure look like the published specifications; the
# container format does not, because Word shows the preview bitmap unless Visio or
# PowerPoint is installed.
NOTATION_FONT = 'Arial'


def _sub_element(parent, tag, **attrs):
    el = parent.makeelement('{%s}%s' % (_A, tag), {k: str(v) for k, v in attrs.items()})
    parent.append(el)
    return el


def _apply_node_style(shape, node):
    """Give a type box the drop shadow the OPC UA notation uses.

    PowerPoint has no 'type box' primitive, so the shadow is an effect on the shape
    itself. It has to be real DrawingML rather than a second grey rectangle, or an editor
    moving the box would leave the shadow behind.
    """
    if not node.is_type:
        return
    sp_pr = shape.fill._xPr
    for existing in sp_pr.findall('{%s}effectLst' % _A):
        sp_pr.remove(existing)
    effect = _sub_element(sp_pr, 'effectLst')
    shadow = _sub_element(effect, 'outerShdw', blurRad=0, dist=38100, dir=2700000,
                          algn='tl', rotWithShape=0)
    _sub_element(shadow, 'srgbClr', val=TYPE_SHADOW)


def _apply_edge_head(connector, edge):
    """Distinguish the reference types the way the notation does.

    HasTypeDefinition is drawn with a large solid head, hierarchical references with the
    ordinary open one. PowerPoint has no double-arrowhead line end, so the distinction is
    carried by head shape and size and stated in the notation legend, rather than being
    approximated with a free-floating second triangle that an editor would detach.
    """
    ln = connector.line._get_or_add_ln()
    for existing in ln.findall('{%s}tailEnd' % _A):
        ln.remove(existing)
    if edge.thick:
        _sub_element(ln, 'tailEnd', type='triangle', w='lg', len='lg')
    else:
        _sub_element(ln, 'tailEnd', type='stealth', w='med', len='med')


def write_pptx(diagram, path):
    if isinstance(diagram, Sequence):
        size = write_sequence_pptx(diagram, path)
    else:
        size = write_flowchart_pptx(diagram, path)
    normalize_pptx(path)
    return size


def normalize_pptx(path):
    """Strip the timestamps python-pptx writes, so the build is byte-reproducible.

    A .pptx carries its creation and modification time in docProps/core.xml and in every
    ZIP entry header. Left alone they change on every run, which would make the embedded
    figure — and therefore the .docx — differ from build to build and destroy the value
    of a clean git diff.
    """
    import re as _re
    import zipfile as _zipfile

    with _zipfile.ZipFile(path) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}
    core = data.get('docProps/core.xml')
    if core is not None:
        text = core.decode('utf-8')
        text = _re.sub(r'>[^<]*</dcterms:created>', '>2026-01-01T00:00:00Z</dcterms:created>',
                       text)
        text = _re.sub(r'>[^<]*</dcterms:modified>',
                       '>2026-01-01T00:00:00Z</dcterms:modified>', text)
        text = _re.sub(r'>[^<]*</cp:lastModifiedBy>', '></cp:lastModifiedBy>', text)
        text = _re.sub(r'>[^<]*</cp:revision>', '>1</cp:revision>', text)
        data['docProps/core.xml'] = text.encode('utf-8')
    with _zipfile.ZipFile(path, 'w', _zipfile.ZIP_DEFLATED) as z:
        for n in names:
            info = _zipfile.ZipInfo(n, date_time=(2026, 1, 1, 0, 0, 0))
            info.compress_type = _zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            z.writestr(info, data[n])


def write_flowchart_pptx(g, path):
    w, h = layout(g)
    prs = Presentation()
    prs.slide_width = Emu(w * EMU_PER_PX)
    prs.slide_height = Emu(h * EMU_PER_PX)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for sid, (x0, y0, x1, y1) in subgraph_bounds(g).items():
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Emu(int(x0 * EMU_PER_PX)), Emu(int(y0 * EMU_PER_PX)),
                                     Emu(int((x1 - x0) * EMU_PER_PX)),
                                     Emu(int((y1 - y0) * EMU_PER_PX)))
        box.fill.solid()
        box.fill.fore_color.rgb = SUBGRAPH_FILL
        box.line.color.rgb = BORDER
        box.line.width = Pt(0.75)
        tf = box.text_frame
        tf.text = g.subgraphs.get(sid, sid)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = TEXT

    notation = g.opcua
    shapes = {}
    for nid in g.order:
        n = g.nodes[nid]
        shape_kind = {'round': MSO_SHAPE.ROUNDED_RECTANGLE,
                      'circle': MSO_SHAPE.OVAL,
                      'diamond': MSO_SHAPE.DIAMOND,
                      # `subroutine` was parsed but missing here, so `[[ ]]` drew the same
                      # rectangle as `[ ]` and a type was indistinguishable from a Node.
                      'subroutine': MSO_SHAPE.RECTANGLE}.get(n.shape, MSO_SHAPE.RECTANGLE)
        sp = slide.shapes.add_shape(shape_kind,
                                    Emu(int(n.x * EMU_PER_PX)), Emu(int(n.y * EMU_PER_PX)),
                                    Emu(int(n.w * EMU_PER_PX)), Emu(int(n.h * EMU_PER_PX)))
        sp.fill.solid()
        sp.fill.fore_color.rgb = FILL
        sp.line.color.rgb = BORDER
        sp.line.width = Pt(1)
        _apply_node_style(sp, n)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.text = n.label
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(9)
                r.font.color.rgb = TEXT
                if notation:
                    r.font.name = NOTATION_FONT
                # Only touched when it carries meaning. Assigning False writes an
                # explicit i="0" into every run, which changed the bytes of diagrams
                # this feature does not apply to at all.
                if n.is_abstract:
                    r.font.italic = True
        shapes[nid] = sp

    for e in g.edges:
        a, b = g.nodes[e.src], g.nodes[e.dst]
        conn = slide.shapes.add_connector(
            # The published figures route with elbows, not diagonals. It is also what
            # keeps a wide figure readable: straight lines between distant boxes cross
            # every row between them.
            MSO_CONNECTOR.ELBOW if notation else MSO_CONNECTOR.STRAIGHT,
            Emu(int((a.x + a.w / 2) * EMU_PER_PX)), Emu(int((a.y + a.h) * EMU_PER_PX)),
            Emu(int((b.x + b.w / 2) * EMU_PER_PX)), Emu(int(b.y * EMU_PER_PX)))
        conn.line.color.rgb = BORDER
        conn.line.width = Pt(1)
        if e.dashed:
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if notation:
            _apply_edge_head(conn, e)
        try:
            conn.begin_connect(shapes[e.src], 2)
            conn.end_connect(shapes[e.dst], 0)
        except (KeyError, ValueError):
            pass
        if e.label:
            # Guideline 1 gets an editable object, so the labels have to live in the
            # PowerPoint too; drawing them only in the preview bitmap would hand an
            # editor a diagram whose edges say nothing.
            tb = slide.shapes.add_textbox(
                Emu(int(e.lx * EMU_PER_PX)), Emu(int(e.ly * EMU_PER_PX)),
                Emu(int((e.lw + 8) * EMU_PER_PX)), Emu(int(LABEL_H * EMU_PER_PX)))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.text = e.label
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(8)
                r.font.color.rgb = BORDER

    prs.save(path)
    return w, h


# --------------------------------------------------------------------------- preview


def write_preview(diagram, path, scale=2):
    """A PNG matching the PowerPoint layout, used as the OLE preview image."""
    if isinstance(diagram, Sequence):
        return write_sequence_preview(diagram, path, scale)
    return write_flowchart_preview(diagram, path, scale)


def write_flowchart_preview(g, path, scale=2):
    w, h = layout(g)
    img = Image.new('RGB', (w * scale, h * scale), 'white')
    d = ImageDraw.Draw(img)
    font = _font(9 * scale)
    title_font = _font(10 * scale, bold=True)
    notation = g.opcua

    for sid, (x0, y0, x1, y1) in subgraph_bounds(g).items():
        d.rounded_rectangle([x0 * scale, y0 * scale, x1 * scale, y1 * scale],
                            radius=6 * scale, outline=(0x44, 0x44, 0x44),
                            fill=(0xF7, 0xF7, 0xF7), width=max(1, scale // 2))
        d.text(((x0 + 8) * scale, (y0 + 5) * scale), g.subgraphs.get(sid, sid),
               fill=(0x11, 0x11, 0x11), font=title_font)

    for e in g.edges:
        a, b = g.nodes[e.src], g.nodes[e.dst]
        x1, y1 = (a.x + a.w / 2) * scale, (a.y + a.h) * scale
        x2, y2 = (b.x + b.w / 2) * scale, b.y * scale
        if notation and abs(x2 - x1) > 1:
            # The same elbow the PowerPoint draws: down to the midpoint between the rows,
            # across, then down into the target. Drawn here too because the preview is
            # what a reviewer looks at.
            mid = (y1 + y2) / 2
            _line(d, x1, y1, x1, mid, scale, e.dashed)
            _line(d, x1, mid, x2, mid, scale, e.dashed)
            _line(d, x2, mid, x2, y2, scale, e.dashed)
            _arrow(d, x2, mid, x2, y2, scale, e.thick)
        else:
            _line(d, x1, y1, x2, y2, scale, e.dashed)
            _arrow(d, x1, y1, x2, y2, scale, notation and e.thick)
        if e.label:
            d.text((e.lx * scale, e.ly * scale), e.label,
                   fill=(0x44, 0x44, 0x44), font=_font(8 * scale))

    for nid in g.order:
        n = g.nodes[nid]
        box = [n.x * scale, n.y * scale, (n.x + n.w) * scale, (n.y + n.h) * scale]
        outline, fill = (0x44, 0x44, 0x44), (0xEE, 0xF3, 0xFA)
        width = max(1, scale // 2)
        if not notation:
            d.rounded_rectangle(box, radius=3 * scale, outline=outline,
                                fill=fill, width=width)
            _centred_text(d, n, font, scale)
            continue
        if n.is_type:
            # The grey offset block behind a type box, matching the drop shadow the
            # PowerPoint carries. Drawn first so the box sits on it.
            off = 4 * scale
            d.rectangle([box[0] + off, box[1] + off, box[2] + off, box[3] + off],
                        fill=(0xBF, 0xBF, 0xBF))
        if n.shape == 'round':
            d.rounded_rectangle(box, radius=min(n.h, 22) * scale // 2, outline=outline,
                                fill=fill, width=width)
        elif n.shape == 'circle':
            d.ellipse(box, outline=outline, fill=fill, width=width)
        elif n.shape == 'diamond':
            cx, cy = (box[0] + box[2]) / 2, (box[1] + box[3]) / 2
            d.polygon([(cx, box[1]), (box[2], cy), (cx, box[3]), (box[0], cy)],
                      outline=outline, fill=fill)
        else:
            d.rectangle(box, outline=outline, fill=fill, width=width)
        _centred_text(d, n, _font(9 * scale, italic=n.is_abstract, arial=True), scale)

    img.save(path, 'PNG')
    return img.size


def _line(d, x1, y1, x2, y2, scale, dashed):
    """A straight edge; dashed when the source marked it `-.->`.

    The dash is semantic in these diagrams — it separates a derived or optional relation
    from a structural one — so a preview that drew every edge solid would misreport the
    model.
    """
    width = max(1, scale // 2)
    if not dashed:
        d.line([x1, y1, x2, y2], fill=(0x44, 0x44, 0x44), width=width)
        return
    import math
    total = math.hypot(x2 - x1, y2 - y1)
    if total == 0:
        return
    dash, gap = 6 * scale, 4 * scale
    pos = 0.0
    while pos < total:
        end = min(pos + dash, total)
        d.line([x1 + (x2 - x1) * pos / total, y1 + (y2 - y1) * pos / total,
                x1 + (x2 - x1) * end / total, y1 + (y2 - y1) * end / total],
               fill=(0x44, 0x44, 0x44), width=width)
        pos = end + gap


def _arrow(d, x1, y1, x2, y2, scale, thick=False):
    import math
    ang = math.atan2(y2 - y1, x2 - x1)
    size = (9 if thick else 6) * scale
    if thick:
        # HasTypeDefinition gets the large solid head the notation uses, so the preview
        # reports the same reference kind as the embedded object.
        p1 = (x2 + size * math.cos(ang + 2.6), y2 + size * math.sin(ang + 2.6))
        p2 = (x2 + size * math.cos(ang - 2.6), y2 + size * math.sin(ang - 2.6))
        d.polygon([(x2, y2), p1, p2], fill=(0x44, 0x44, 0x44))
        return
    for delta in (2.6, -2.6):
        d.line([x2, y2,
                x2 + size * math.cos(ang + delta),
                y2 + size * math.sin(ang + delta)],
               fill=(0x44, 0x44, 0x44), width=max(1, scale // 2))


def _centred_text(d, n, font, scale):
    lines = _wrap(d, n.label, font, (n.w - 12) * scale)
    line_h = (font.size + 3 * scale) if hasattr(font, 'size') else 12 * scale
    total = len(lines) * line_h
    y = (n.y + n.h / 2) * scale - total / 2
    for line in lines:
        try:
            tw = d.textlength(line, font=font)
        except AttributeError:
            tw = len(line) * 6 * scale
        d.text(((n.x + n.w / 2) * scale - tw / 2, y), line,
               fill=(0x11, 0x11, 0x11), font=font)
        y += line_h


def _wrap(d, text, font, max_px):
    out = []
    for hard in text.split('\n'):
        words = hard.split()
        line = ''
        for w in words:
            probe = (line + ' ' + w).strip()
            try:
                width = d.textlength(probe, font=font)
            except AttributeError:
                width = len(probe) * 6
            if width > max_px and line:
                out.append(line)
                line = w
            else:
                line = probe
        out.append(line)
    return [x for x in out if x]


def _font(size, bold=False, italic=False, arial=False):
    if arial:
        # The notation figures use Arial to match the published specifications, so the
        # preview has to as well or it disagrees with the object it previews.
        names = (('arialbd.ttf',) if bold else
                 ('ariali.ttf',) if italic else ('arial.ttf',)) + ('segoeui.ttf',)
    elif bold:
        names = ('seguisb.ttf', 'segoeuib.ttf')
    elif italic:
        # An abstract type is italic in the OPC UA notation, so the preview needs the
        # italic face or it disagrees with the embedded PowerPoint.
        names = ('segoeuii.ttf', 'ariali.ttf', 'segoeui.ttf')
    else:
        names = ('segoeui.ttf', 'arial.ttf')
    for name in names:
        try:
            f = ImageFont.truetype(name, int(size))
            f.size = int(size)
            return f
        except OSError:
            continue
    f = ImageFont.load_default()
    try:
        f.size = int(size)
    except AttributeError:
        pass
    return f


# --------------------------------------------------------------------------- sequence

PARTICIPANT_RE = re.compile(r'^(?:participant|actor)\s+(?P<id>[A-Za-z0-9_]+)'
                            r'(?:\s+as\s+(?P<label>.+))?$')
MESSAGE_RE = re.compile(r'^(?P<src>[A-Za-z0-9_]+)\s*'
                        r'(?P<arrow>-->>|->>|-->|->|-x|--x)\s*'
                        r'(?P<dst>[A-Za-z0-9_]+)\s*:\s*(?P<label>.*)$')
NOTE_RE = re.compile(r'^Note\s+(?:over|left of|right of)\s+(?P<who>[^:]+):\s*(?P<text>.*)$',
                     re.IGNORECASE)
BLOCK_RE = re.compile(r'^(?P<kind>loop|alt|else|opt|par|critical|rect)\b\s*(?P<label>.*)$')

PARTICIPANT_W = 190
PARTICIPANT_H = 44
PARTICIPANT_GAP = 34
ROW_H = 40
SEQ_MARGIN = 26
BAND_H = 30


class Sequence:
    def __init__(self):
        self.participants = []
        self.labels = {}
        self.rows = []     # ('msg', src, dst, label, dashed) | ('note', text) | ('band', text)

    def participant(self, pid, label=None):
        if pid not in self.labels:
            self.participants.append(pid)
            self.labels[pid] = label or pid
        elif label:
            self.labels[pid] = label
        return pid


def parse_sequence(source):
    s = Sequence()
    depth = 0
    for raw in source.splitlines():
        line = raw.strip()
        if not line or line.startswith('%%') or line == 'autonumber':
            continue
        if line.startswith('sequenceDiagram'):
            continue
        if line == 'end':
            depth = max(0, depth - 1)
            continue
        m = PARTICIPANT_RE.match(line)
        if m:
            s.participant(m.group('id'), _clean(m.group('label') or m.group('id')))
            continue
        m = NOTE_RE.match(line)
        if m:
            s.rows.append(('note', _clean(m.group('text'))))
            continue
        m = MESSAGE_RE.match(line)
        if m:
            src = s.participant(m.group('src'))
            dst = s.participant(m.group('dst'))
            dashed = m.group('arrow').startswith('--')
            s.rows.append(('msg', src, dst, _clean(m.group('label')), dashed))
            continue
        m = BLOCK_RE.match(line)
        if m:
            depth += 1
            s.rows.append(('band', '%s %s' % (m.group('kind'), _clean(m.group('label')))))
            continue
        raise ValueError('unsupported mermaid sequence line: %r' % raw)
    return s


def sequence_layout(s):
    n = max(1, len(s.participants))
    width = SEQ_MARGIN * 2 + n * PARTICIPANT_W + (n - 1) * PARTICIPANT_GAP
    height = SEQ_MARGIN * 2 + PARTICIPANT_H + 20 + len(s.rows) * ROW_H + 20
    return int(width), int(height)


def _participant_x(s, pid):
    i = s.participants.index(pid)
    return SEQ_MARGIN + i * (PARTICIPANT_W + PARTICIPANT_GAP) + PARTICIPANT_W / 2


def write_sequence_pptx(s, path):
    w, h = sequence_layout(s)
    prs = Presentation()
    prs.slide_width = Emu(w * EMU_PER_PX)
    prs.slide_height = Emu(h * EMU_PER_PX)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = SEQ_MARGIN
    body_top = top + PARTICIPANT_H + 20

    for i, pid in enumerate(s.participants):
        x = SEQ_MARGIN + i * (PARTICIPANT_W + PARTICIPANT_GAP)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Emu(int(x * EMU_PER_PX)), Emu(int(top * EMU_PER_PX)),
                                     Emu(int(PARTICIPANT_W * EMU_PER_PX)),
                                     Emu(int(PARTICIPANT_H * EMU_PER_PX)))
        box.fill.solid()
        box.fill.fore_color.rgb = FILL
        box.line.color.rgb = BORDER
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = s.labels[pid]
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.color.rgb = TEXT
        cx = x + PARTICIPANT_W / 2
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(int(cx * EMU_PER_PX)), Emu(int(body_top * EMU_PER_PX)),
            Emu(int(cx * EMU_PER_PX)), Emu(int((h - SEQ_MARGIN) * EMU_PER_PX)))
        line.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        line.line.width = Pt(0.75)

    for k, entry in enumerate(s.rows):
        y = body_top + k * ROW_H + ROW_H / 2
        if entry[0] in ('note', 'band'):
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(int(SEQ_MARGIN * EMU_PER_PX)), Emu(int((y - BAND_H / 2) * EMU_PER_PX)),
                Emu(int((w - 2 * SEQ_MARGIN) * EMU_PER_PX)), Emu(int(BAND_H * EMU_PER_PX)))
            band.fill.solid()
            band.fill.fore_color.rgb = SUBGRAPH_FILL
            band.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            band.line.width = Pt(0.75)
            tf = band.text_frame
            tf.text = entry[1]
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for r in para.runs:
                    r.font.size = Pt(8)
                    r.font.italic = True
                    r.font.color.rgb = TEXT
            continue
        _, src, dst, label, dashed = entry
        x1 = _participant_x(s, src)
        x2 = _participant_x(s, dst)
        if src == dst:
            x2 = x1 + 70
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(int(x1 * EMU_PER_PX)), Emu(int(y * EMU_PER_PX)),
            Emu(int(x2 * EMU_PER_PX)), Emu(int(y * EMU_PER_PX)))
        conn.line.color.rgb = BORDER
        conn.line.width = Pt(1)
        if dashed:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        tb = slide.shapes.add_textbox(
            Emu(int(min(x1, x2) * EMU_PER_PX)),
            Emu(int((y - ROW_H / 2 + 2) * EMU_PER_PX)),
            Emu(int(max(abs(x2 - x1), 120) * EMU_PER_PX)),
            Emu(int((ROW_H / 2) * EMU_PER_PX)))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = label
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(8)
                r.font.color.rgb = TEXT

    prs.save(path)
    return w, h


def write_sequence_preview(s, path, scale=2):
    w, h = sequence_layout(s)
    img = Image.new('RGB', (w * scale, h * scale), 'white')
    d = ImageDraw.Draw(img)
    font = _font(9 * scale)
    small = _font(8 * scale)
    top = SEQ_MARGIN
    body_top = top + PARTICIPANT_H + 20

    for i, pid in enumerate(s.participants):
        x = SEQ_MARGIN + i * (PARTICIPANT_W + PARTICIPANT_GAP)
        d.rounded_rectangle([x * scale, top * scale,
                             (x + PARTICIPANT_W) * scale, (top + PARTICIPANT_H) * scale],
                            radius=4 * scale, outline=(0x44, 0x44, 0x44),
                            fill=(0xEE, 0xF3, 0xFA), width=max(1, scale // 2))
        node = Node(pid, s.labels[pid])
        node.x, node.y, node.w, node.h = x, top, PARTICIPANT_W, PARTICIPANT_H
        _centred_text(d, node, font, scale)
        cx = (x + PARTICIPANT_W / 2) * scale
        d.line([cx, body_top * scale, cx, (h - SEQ_MARGIN) * scale],
               fill=(0x99, 0x99, 0x99), width=max(1, scale // 2))

    for k, entry in enumerate(s.rows):
        y = body_top + k * ROW_H + ROW_H / 2
        if entry[0] in ('note', 'band'):
            d.rectangle([SEQ_MARGIN * scale, (y - BAND_H / 2) * scale,
                         (w - SEQ_MARGIN) * scale, (y + BAND_H / 2) * scale],
                        outline=(0xAA, 0xAA, 0xAA), fill=(0xF7, 0xF7, 0xF7),
                        width=max(1, scale // 2))
            band = Node('band', entry[1])
            band.x, band.y = SEQ_MARGIN, y - BAND_H / 2
            band.w, band.h = w - 2 * SEQ_MARGIN, BAND_H
            _centred_text(d, band, small, scale)
            continue
        _, src, dst, label, dashed = entry
        x1 = _participant_x(s, src)
        x2 = _participant_x(s, dst)
        if src == dst:
            x2 = x1 + 70
        d.line([x1 * scale, y * scale, x2 * scale, y * scale],
               fill=(0x44, 0x44, 0x44), width=max(1, scale // 2))
        _arrow(d, x1 * scale, y * scale, x2 * scale, y * scale, scale)
        holder = Node('lbl', label)
        holder.x = min(x1, x2)
        holder.y = y - ROW_H / 2
        holder.w = max(abs(x2 - x1), 120)
        holder.h = ROW_H / 2
        _centred_text(d, holder, small, scale)

    img.save(path, 'PNG')
    return img.size

